"""File readers for extracting text from various document formats."""

import os
from pathlib import Path
from typing import Optional
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Global OCR reader instance (initialized lazily)
_ocr_reader = None


def get_ocr_reader():
    """Get or initialize the OCR reader (singleton pattern)."""
    global _ocr_reader
    
    if _ocr_reader is not None:
        return _ocr_reader
    
    from extractor.config import OCR_BACKEND, OCR_LANGUAGES, OCR_USE_GPU
    
    if OCR_BACKEND == "easyocr":
        try:
            import easyocr
            logger.info(f"Initializing EasyOCR with GPU={OCR_USE_GPU}")
            _ocr_reader = easyocr.Reader(OCR_LANGUAGES, gpu=OCR_USE_GPU)
            logger.info("EasyOCR initialized successfully")
            return _ocr_reader
        except Exception as e:
            logger.error(f"Failed to initialize EasyOCR: {e}")
            return None
    
    elif OCR_BACKEND == "doctr":
        try:
            from doctr.io import DocumentFile
            from doctr.models import ocr_predictor
            logger.info("Initializing doctr OCR")
            _ocr_reader = ocr_predictor(pretrained=True)
            logger.info("doctr OCR initialized successfully")
            return _ocr_reader
        except Exception as e:
            logger.error(f"Failed to initialize doctr: {e}")
            return None
    
    elif OCR_BACKEND == "paddleocr":
        try:
            from paddleocr import PaddleOCR
            logger.info("Initializing PaddleOCR")
            _ocr_reader = PaddleOCR(use_angle_cls=True, lang='en', use_gpu=OCR_USE_GPU)
            logger.info("PaddleOCR initialized successfully")
            return _ocr_reader
        except Exception as e:
            logger.error(f"Failed to initialize PaddleOCR: {e}")
            return None
    
    return None


def read_txt(file_path: Path) -> str:
    """Read text from a TXT file."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        logger.error(f"Error reading TXT {file_path}: {e}")
        return ""


def read_pdf(file_path: Path) -> str:
    """Read text from a PDF file using PyMuPDF."""
    try:
        import fitz  # PyMuPDF
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        
        # If no text extracted (scanned PDF), try OCR
        if len(text.strip()) < 50:
            logger.info(f"PDF appears to be scanned, attempting OCR: {file_path.name}")
            text = read_pdf_with_ocr(file_path)
        
        return text
    except Exception as e:
        logger.error(f"Error reading PDF {file_path}: {e}")
        return ""


def read_pdf_with_ocr(file_path: Path) -> str:
    """Read scanned PDF using OCR."""
    try:
        import fitz
        from PIL import Image
        import io
        
        text_parts = []
        
        with fitz.open(file_path) as doc:
            for page_num, page in enumerate(doc):
                # Render page to image
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                
                # OCR the image
                page_text = _ocr_image(img)
                if page_text:
                    text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
        
        return "\n\n".join(text_parts)
    except Exception as e:
        logger.error(f"Error in PDF OCR for {file_path}: {e}")
        return ""


def read_docx(file_path: Path) -> str:
    """Read text from a DOCX file."""
    try:
        from docx import Document
        doc = Document(file_path)
        text = []
        for para in doc.paragraphs:
            text.append(para.text)
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text.append(cell.text)
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error reading DOCX {file_path}: {e}")
        return ""


def read_pptx(file_path: Path) -> str:
    """Read text from a PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text.append(cell.text)
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error reading PPTX {file_path}: {e}")
        return ""


def read_xlsx(file_path: Path) -> str:
    """Read text from an XLSX file."""
    try:
        import pandas as pd
        dfs = pd.read_excel(file_path, sheet_name=None)
        text = []
        for sheet_name, df in dfs.items():
            text.append(f"Sheet: {sheet_name}")
            text.append(df.to_string())
        return "\n".join(text)
    except Exception as e:
        logger.error(f"Error reading XLSX {file_path}: {e}")
        return ""


def _ocr_image(image) -> str:
    """Perform OCR on a PIL Image using configured backend."""
    from extractor.config import OCR_BACKEND
    
    reader = get_ocr_reader()
    if reader is None:
        logger.warning("OCR reader not available")
        return ""
    
    try:
        if OCR_BACKEND == "easyocr":
            import numpy as np
            # Convert PIL image to numpy array
            img_array = np.array(image)
            results = reader.readtext(img_array)
            # Extract text from results
            text = " ".join([result[1] for result in results])
            return text
        
        elif OCR_BACKEND == "doctr":
            from doctr.io import DocumentFile
            import numpy as np
            import tempfile
            
            # Save image temporarily
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                image.save(tmp.name)
                doc = DocumentFile.from_images(tmp.name)
                result = reader(doc)
                # Extract text
                text = ""
                for page in result.pages:
                    for block in page.blocks:
                        for line in block.lines:
                            text += " ".join([word.value for word in line.words]) + "\n"
                os.unlink(tmp.name)
            return text
        
        elif OCR_BACKEND == "paddleocr":
            import numpy as np
            img_array = np.array(image)
            result = reader.ocr(img_array, cls=True)
            text = ""
            if result and result[0]:
                for line in result[0]:
                    if line[1]:
                        text += line[1][0] + " "
            return text
        
    except Exception as e:
        logger.error(f"OCR error: {e}")
        return ""
    
    return ""


def read_image(file_path: Path) -> str:
    """Extract text from an image using OCR."""
    try:
        from PIL import Image
        
        img = Image.open(file_path)
        
        # Convert to RGB if necessary
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        text = _ocr_image(img)
        
        if text:
            logger.info(f"OCR extracted {len(text)} chars from {file_path.name}")
        else:
            logger.warning(f"OCR extracted no text from {file_path.name}")
        
        return text
    except Exception as e:
        logger.error(f"Error reading image {file_path}: {e}")
        return ""


def read_file(file_path: Path) -> str:
    """Read text from a file based on its extension."""
    suffix = file_path.suffix.lower()
    
    readers = {
        '.txt': read_txt,
        '.pdf': read_pdf,
        '.docx': read_docx,
        '.doc': read_docx,
        '.pptx': read_pptx,
        '.ppt': read_pptx,
        '.xlsx': read_xlsx,
        '.xls': read_xlsx,
        '.png': read_image,
        '.jpg': read_image,
        '.jpeg': read_image,
    }
    
    reader = readers.get(suffix)
    if reader:
        return reader(file_path)
    else:
        logger.warning(f"Unsupported file format: {suffix} for {file_path}")
        return ""
